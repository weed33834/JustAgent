"""Document model, parsing, and lifecycle management for the knowledge base.

Supports ingestion of multiple document formats (PDF, Word, Excel, PPT,
Markdown, HTML, plain text) with format detection based on file extension
and content sniffing. Document parsing uses the standard library for text
and Markdown; binary formats (PDF, Word, Excel, PPT) use lazy imports with
graceful fallback so the module remains importable when those optional
dependencies are not installed.

The :class:`DocumentLifecycleManager` provides versioning and archival:
each edit creates a new version snapshot, and archived documents are
retained but excluded from active search.

Design:

* :class:`DocumentType` — enum of supported formats.
* :class:`DocumentStatus` — lifecycle state (active / archived / deleted).
* :class:`Chunk` — a text fragment extracted from a document, ready for
  vector indexing.
* :class:`DocumentVersion` — an immutable snapshot of a document's content
  at a point in time.
* :class:`Document` — the central Pydantic model tying everything together.
* :class:`TextChunker` — splits raw text into overlapping chunks.
* :class:`DocumentParser` — detects format, reads the file, and returns a
  :class:`Document`.
* :class:`DocumentLifecycleManager` — tracks versions, archives, restores,
  and deletes documents.
"""

from __future__ import annotations

import hashlib
import logging
import re
import time
import uuid
from collections.abc import Iterator
from datetime import UTC, datetime
from enum import Enum
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from pydantic import BaseModel, Field

from justagent.utils import utcnow  # noqa: F401 - re-exported for backwards compatibility

logger = logging.getLogger("justagent.knowledge")


# ---------------------------------------------------------------------------
# Enums
# ---------------------------------------------------------------------------


class DocumentType(str, Enum):  # noqa: UP042
    """Supported document formats."""

    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    PPT = "ppt"
    MARKDOWN = "markdown"
    HTML = "html"
    PLAIN_TEXT = "plain_text"
    UNKNOWN = "unknown"


class DocumentStatus(str, Enum):  # noqa: UP042
    """Lifecycle state of a document.

    * ``active`` — searchable and included in RAG results.
    * ``archived`` — retained on disk but excluded from active search.
    * ``deleted`` — soft-deleted; may be purged by a background sweep.
    """

    ACTIVE = "active"
    ARCHIVED = "archived"
    DELETED = "deleted"


# ---------------------------------------------------------------------------
# Extension / type maps
# ---------------------------------------------------------------------------

_EXTENSION_MAP: dict[str, DocumentType] = {
    ".pdf": DocumentType.PDF,
    ".doc": DocumentType.WORD,
    ".docx": DocumentType.WORD,
    ".xls": DocumentType.EXCEL,
    ".xlsx": DocumentType.EXCEL,
    ".ppt": DocumentType.PPT,
    ".pptx": DocumentType.PPT,
    ".md": DocumentType.MARKDOWN,
    ".markdown": DocumentType.MARKDOWN,
    ".html": DocumentType.HTML,
    ".htm": DocumentType.HTML,
    ".xhtml": DocumentType.HTML,
    ".txt": DocumentType.PLAIN_TEXT,
    ".text": DocumentType.PLAIN_TEXT,
    ".log": DocumentType.PLAIN_TEXT,
    ".csv": DocumentType.PLAIN_TEXT,
    ".json": DocumentType.PLAIN_TEXT,
    ".yaml": DocumentType.PLAIN_TEXT,
    ".yml": DocumentType.PLAIN_TEXT,
    ".py": DocumentType.PLAIN_TEXT,
    ".js": DocumentType.PLAIN_TEXT,
    ".ts": DocumentType.PLAIN_TEXT,
    ".rs": DocumentType.PLAIN_TEXT,
    ".go": DocumentType.PLAIN_TEXT,
    ".java": DocumentType.PLAIN_TEXT,
    ".c": DocumentType.PLAIN_TEXT,
    ".cpp": DocumentType.PLAIN_TEXT,
    ".h": DocumentType.PLAIN_TEXT,
    ".sh": DocumentType.PLAIN_TEXT,
}


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------


class Chunk(BaseModel):
    """A text chunk extracted from a document for vector indexing.

    Attributes:
        id: Unique chunk identifier.
        document_id: ID of the parent document.
        content: The chunk text.
        index: Position of this chunk within the parent document.
        metadata: Arbitrary key-value metadata (e.g. heading, page number).
        token_count: Approximate token count, computed from content.
    """

    id: str = Field(default_factory=lambda: uuid.uuid4().hex)
    document_id: str
    content: str
    index: int = 0
    metadata: dict[str, Any] = Field(default_factory=dict)
    token_count: int = 0

    def model_post_init(self, __context: Any) -> None:
        """Compute ``token_count`` if it was not explicitly set."""
        if self.token_count == 0 and self.content:
            self.token_count = _estimate_tokens(self.content)


class DocumentVersion(BaseModel):
    """An immutable snapshot of a document's content at a point in time.

    Attributes:
        version: Monotonically increasing version number.
        content: Full text content at this version.
        content_hash: SHA-256 hex digest of ``content``.
        created_at: Unix timestamp of when the version was captured.
        metadata: Snapshot of document metadata at this version.
    """

    version: int
    content: str
    content_hash: str
    created_at: float = Field(default_factory=lambda: time.time())
    metadata: dict[str, Any] = Field(default_factory=dict)


class Document(BaseModel):
    """A knowledge-base document with full lifecycle metadata.

    Attributes:
        id: Unique document identifier.
        title: Human-readable title.
        source: Original file path or URL.
        type: Detected :class:`DocumentType`.
        content: Full extracted text content.
        chunks: List of :class:`Chunk` objects derived from ``content``.
        metadata: Arbitrary key-value metadata.
        status: Current lifecycle status.
        version: Current version number (starts at 1).
        versions: History of :class:`DocumentVersion` snapshots.
        content_hash: SHA-256 hex digest of ``content``.
        created_at: Unix timestamp of creation.
        updated_at: Unix timestamp of last modification.
    """

    id: str = Field(default_factory=lambda: uuid.uuid4().hex)
    title: str
    source: str = ""
    type: DocumentType = DocumentType.PLAIN_TEXT
    content: str = ""
    chunks: list[Chunk] = Field(default_factory=list)
    metadata: dict[str, Any] = Field(default_factory=dict)
    status: DocumentStatus = DocumentStatus.ACTIVE
    version: int = 1
    versions: list[DocumentVersion] = Field(default_factory=list)
    content_hash: str = ""
    created_at: float = Field(default_factory=lambda: time.time())
    updated_at: float = Field(default_factory=lambda: time.time())

    def model_post_init(self, __context: Any) -> None:
        """Compute ``content_hash`` if not already set."""
        if not self.content_hash and self.content:
            self.content_hash = _sha256_text(self.content)

    # ------------------------------------------------------------------
    # Convenience properties
    # ------------------------------------------------------------------

    @property
    def is_active(self) -> bool:
        """True if the document is active and searchable."""
        return self.status is DocumentStatus.ACTIVE

    @property
    def is_archived(self) -> bool:
        """True if the document is archived."""
        return self.status is DocumentStatus.ARCHIVED

    @property
    def is_deleted(self) -> bool:
        """True if the document is soft-deleted."""
        return self.status is DocumentStatus.DELETED

    @property
    def token_count(self) -> int:
        """Approximate total token count across all chunks."""
        return sum(c.token_count for c in self.chunks)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _sha256_text(text: str) -> str:
    """Return the SHA-256 hex digest of ``text``."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _estimate_tokens(text: str) -> int:
    """Estimate token count using the ~4 chars-per-token heuristic.

    This is a coarse approximation suitable for budgeting and logging.
    For precise counts, use a proper tokenizer at the call site.
    """
    return max(1, len(text) // 4)


def detect_type(path: Path | str) -> DocumentType:
    """Detect :class:`DocumentType` from a file path extension.

    Falls back to :attr:`DocumentType.UNKNOWN` for unrecognised extensions.
    """
    ext = Path(path).suffix.lower()
    return _EXTENSION_MAP.get(ext, DocumentType.UNKNOWN)


def format_timestamp(ts: float) -> str:
    """Format a Unix timestamp as an ISO-8601 UTC string."""
    return datetime.fromtimestamp(ts, tz=UTC).isoformat()


# ---------------------------------------------------------------------------
# Text chunking
# ---------------------------------------------------------------------------


class TextChunker:
    """Split raw text into overlapping chunks for vector indexing.

    The chunker first splits on paragraph boundaries (double newlines).
    Paragraphs longer than ``chunk_size`` characters are further split on
    sentence boundaries. Each chunk overlaps the previous one by
    ``chunk_overlap`` characters to preserve context across boundaries.

    Example::

        >>> chunker = TextChunker(chunk_size=200, chunk_overlap=50)
        >>> chunks = chunker.chunk("Hello world. " * 100)
        >>> len(chunks) > 1
        True
    """

    # Sentence-end pattern: period, exclamation, or question mark followed
    # by whitespace and a capital letter or end of string.
    _SENTENCE_END = re.compile(r"(?<=[.!?])\s+(?=[A-Z])")
    # Whitespace runs to collapse.
    _WHITESPACE = re.compile(r"[ \t]+")

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 50,
    ) -> None:
        if chunk_size <= 0:
            raise ValueError("chunk_size must be positive")
        if chunk_overlap < 0:
            raise ValueError("chunk_overlap must be non-negative")
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap must be smaller than chunk_size")
        if min_chunk_size < 1:
            raise ValueError("min_chunk_size must be at least 1")
        self._chunk_size = chunk_size
        self._chunk_overlap = chunk_overlap
        self._min_chunk_size = min_chunk_size

    @property
    def chunk_size(self) -> int:
        return self._chunk_size

    @property
    def chunk_overlap(self) -> int:
        return self._chunk_overlap

    def chunk(self, text: str) -> list[str]:
        """Split ``text`` into overlapping chunks.

        Returns a list of non-empty chunk strings. Very short inputs
        (below ``min_chunk_size``) are returned as a single chunk.
        """
        text = self._normalize(text)
        if not text.strip():
            return []

        # Start with paragraph-level splits.
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

        # Further split long paragraphs into sentences.
        sentences: list[str] = []
        for para in paragraphs:
            if len(para) <= self._chunk_size:
                sentences.append(para)
            else:
                parts = self._SENTENCE_END.split(para)
                for part in parts:
                    part = part.strip()
                    if part:
                        sentences.append(part)

        # Merge sentences into chunks respecting chunk_size with overlap.
        chunks: list[str] = []
        current = ""
        for sentence in sentences:
            candidate = f"{current} {sentence}".strip() if current else sentence
            if len(candidate) <= self._chunk_size:
                current = candidate
            else:
                if current:
                    chunks.append(current)
                # Start the next chunk with overlap from the tail of the
                # previous chunk (if it was long enough to warrant overlap).
                if self._chunk_overlap > 0 and len(current) > self._min_chunk_size:
                    overlap_text = current[-self._chunk_overlap:]
                    current = f"{overlap_text} {sentence}".strip()
                    # If the single sentence itself exceeds chunk_size,
                    # hard-split it.
                    if len(current) > self._chunk_size:
                        for piece in self._hard_split(current):
                            chunks.append(piece)
                        current = ""
                else:
                    current = sentence
                    # Hard-split if a single sentence is too long.
                    if len(current) > self._chunk_size:
                        for piece in self._hard_split(current):
                            chunks.append(piece)
                        current = ""
        if current:
            chunks.append(current)

        return [c for c in chunks if len(c) >= 1]

    def chunk_document(self, document: Document) -> list[Chunk]:
        """Chunk a :class:`Document`'s content and return :class:`Chunk` list."""
        texts = self.chunk(document.content)
        return [
            Chunk(
                document_id=document.id,
                content=text,
                index=i,
                metadata={"source": document.source, "title": document.title},
            )
            for i, text in enumerate(texts)
        ]

    # ------------------------------------------------------------------
    # Internal
    # ------------------------------------------------------------------

    def _normalize(self, text: str) -> str:
        """Collapse excessive whitespace while preserving paragraph breaks."""
        lines = text.split("\n")
        lines = [self._WHITESPACE.sub(" ", line).rstrip() for line in lines]
        # Re-join and collapse 3+ newlines into 2 (paragraph break).
        text = "\n".join(lines)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _hard_split(self, text: str) -> list[str]:
        """Hard-split ``text`` into pieces no longer than ``chunk_size``."""
        pieces: list[str] = []
        start = 0
        while start < len(text):
            end = min(start + self._chunk_size, len(text))
            # Try to break on a space boundary.
            if end < len(text):
                space = text.rfind(" ", start, end)
                if space > start + self._min_chunk_size:
                    end = space
            piece = text[start:end].strip()
            if piece:
                pieces.append(piece)
            start = end
        return pieces


# ---------------------------------------------------------------------------
# Markdown and HTML parsing (standard library only)
# ---------------------------------------------------------------------------


class _MarkdownStripper:
    """Convert Markdown to plain text using regex substitutions.

    This is not a full Markdown parser — it handles the common syntax
    (headers, bold, italic, links, images, code blocks, lists, quotes,
    horizontal rules) well enough for text extraction and indexing.
    """

    _PATTERNS: list[tuple[re.Pattern[str], str]] = [
        # Fenced code blocks — keep the code, remove fences.
        (re.compile(r"```[^\n]*\n(.*?)```", re.DOTALL), r"\1"),
        # Inline code — keep the code.
        (re.compile(r"`([^`]+)`"), r"\1"),
        # Images ![alt](url) — keep alt text.
        (re.compile(r"!\[([^\]]*)\]\([^)]+\)"), r"\1"),
        # Links [text](url) — keep text.
        (re.compile(r"\[([^\]]*)\]\([^)]+\)"), r"\1"),
        # Reference links [text][ref] — keep text.
        (re.compile(r"\[([^\]]*)\]\[[^\]]*\]"), r"\1"),
        # Bold + italic.
        (re.compile(r"\*\*\*(.+?)\*\*\*"), r"\1"),
        (re.compile(r"___(.+?)___"), r"\1"),
        # Bold.
        (re.compile(r"\*\*(.+?)\*\*"), r"\1"),
        (re.compile(r"__(.+?)__"), r"\1"),
        # Italic.
        (re.compile(r"\*(.+?)\*"), r"\1"),
        (re.compile(r"_(.+?)_"), r"\1"),
        # Strikethrough.
        (re.compile(r"~~(.+?)~~"), r"\1"),
        # Headers — keep text, remove leading #'s.
        (re.compile(r"^#{1,6}\s+", re.MULTILINE), ""),
        # Horizontal rules.
        (re.compile(r"^[-*_]{3,}\s*$", re.MULTILINE), ""),
        # Blockquotes — remove leading >.
        (re.compile(r"^>\s?", re.MULTILINE), ""),
        # Unordered list markers.
        (re.compile(r"^[\s]*[-*+]\s+", re.MULTILINE), ""),
        # Ordered list markers.
        (re.compile(r"^[\s]*\d+\.\s+", re.MULTILINE), ""),
        # Reference link definitions [ref]: url.
        (re.compile(r"^\[[^\]]+\]:\s*\S+.*$", re.MULTILINE), ""),
    ]

    @classmethod
    def strip(cls, markdown: str) -> str:
        """Convert ``markdown`` text to plain text."""
        text = markdown
        for pattern, replacement in cls._PATTERNS:
            text = pattern.sub(replacement, text)
        # Collapse excessive whitespace.
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


class _HTMLTextExtractor(HTMLParser):
    """Extract visible text from HTML using the standard library parser.

    Skips ``<script>`` and ``<style>`` content. Inserts newlines for
    block-level elements to preserve some structure.
    """

    _SKIP_TAGS = frozenset({"script", "style", "noscript", "head"})
    _BLOCK_TAGS = frozenset({
        "p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6",
        "section", "article", "header", "footer", "blockquote", "pre",
        "ul", "ol", "table", "hr",
    })

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        elif tag in self._BLOCK_TAGS:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth = max(0, self._skip_depth - 1)
        elif tag in self._BLOCK_TAGS:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0:
            self._parts.append(data)

    def get_text(self) -> str:
        text = "".join(self._parts)
        # Collapse whitespace.
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


def parse_html(html: str) -> str:
    """Extract plain text from an HTML string."""
    extractor = _HTMLTextExtractor()
    extractor.feed(html)
    return extractor.get_text()


def parse_markdown(md: str) -> str:
    """Convert Markdown to plain text."""
    return _MarkdownStripper.strip(md)


# ---------------------------------------------------------------------------
# Binary format parsers (lazy imports with graceful fallback)
# ---------------------------------------------------------------------------


def _parse_pdf(data: bytes, source: str) -> str:
    """Extract text from PDF bytes using PyPDF2 or pdfplumber (lazy import).

    Falls back to a placeholder message if neither library is installed.
    """
    try:
        import io

        from PyPDF2 import PdfReader
    except ImportError:
        pass
    else:
        try:
            reader = PdfReader(io.BytesIO(data))
            parts: list[str] = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text.strip())
            if parts:
                return "\n\n".join(parts)
        except Exception as exc:
            logger.warning("PyPDF2 failed to parse %s: %s", source, exc)

    try:
        import io

        import pdfplumber
    except ImportError:
        pass
    else:
        try:
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                parts = []
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        parts.append(text.strip())
                if parts:
                    return "\n\n".join(parts)
        except Exception as exc:
            logger.warning("pdfplumber failed to parse %s: %s", source, exc)

    logger.warning(
        "No PDF parser available (install PyPDF2 or pdfplumber) for %s", source
    )
    return f"[PDF content not extracted: {source}. Install PyPDF2 or pdfplumber.]\n"


def _parse_word(data: bytes, source: str) -> str:
    """Extract text from .docx bytes using python-docx (lazy import).

    Falls back to a placeholder message if python-docx is not installed.
    """
    try:
        import io

        from docx import Document as DocxDocument
    except ImportError:
        logger.warning(
            "python-docx not installed; cannot parse Word document %s", source
        )
        return f"[Word content not extracted: {source}. Install python-docx.]\n"

    try:
        doc = DocxDocument(io.BytesIO(data))
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        # Also extract tables.
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)
    except Exception as exc:
        logger.warning("python-docx failed to parse %s: %s", source, exc)
        return f"[Word content not extracted: {source}. Error: {exc}]\n"


def _parse_excel(data: bytes, source: str) -> str:
    """Extract text from .xlsx bytes using openpyxl (lazy import).

    Falls back to a placeholder message if openpyxl is not installed.
    """
    try:
        import io

        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed; cannot parse Excel file %s", source)
        return f"[Excel content not extracted: {source}. Install openpyxl.]\n"

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in wb.worksheets:
            parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [
                    str(cell).strip()
                    for cell in row
                    if cell is not None and str(cell).strip()
                ]
                if cells:
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("openpyxl failed to parse %s: %s", source, exc)
        return f"[Excel content not extracted: {source}. Error: {exc}]\n"


def _parse_ppt(data: bytes, source: str) -> str:
    """Extract text from .pptx bytes using python-pptx (lazy import).

    Falls back to a placeholder message if python-pptx is not installed.
    """
    try:
        import io

        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; cannot parse PPT file %s", source)
        return f"[PPT content not extracted: {source}. Install python-pptx.]\n"

    try:
        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        ]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"## Slide {i}\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as exc:
        logger.warning("python-pptx failed to parse %s: %s", source, exc)
        return f"[PPT content not extracted: {source}. Error: {exc}]\n"


# ---------------------------------------------------------------------------
# Document parser
# ---------------------------------------------------------------------------


class DocumentParser:
    """Parse files and raw text into :class:`Document` objects.

    Format detection is based on file extension (see
    :func:`detect_type`). For binary formats (PDF, Word, Excel, PPT),
    the appropriate library is lazily imported; if the library is not
    installed, a placeholder message is stored as content so the document
    is still created and tracked.

    Example::

        >>> parser = DocumentParser()
        >>> doc = parser.parse_file(Path("README.md"))
        >>> doc.type
        <DocumentType.MARKDOWN: 'markdown'>
    """

    def __init__(self, chunker: TextChunker | None = None) -> None:
        self._chunker = chunker or TextChunker()

    @property
    def chunker(self) -> TextChunker:
        return self._chunker

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def parse_file(
        self,
        path: Path | str,
        *,
        title: str | None = None,
        metadata: dict[str, Any] | None = None,
        auto_chunk: bool = True,
    ) -> Document:
        """Parse a file into a :class:`Document`.

        Args:
            path: Path to the file to parse.
            title: Optional title override. Defaults to the file name.
            metadata: Optional metadata to attach.
            auto_chunk: If True (default), the content is chunked and
                the chunks are attached to the document.

        Raises:
            FileNotFoundError: If ``path`` does not exist.
        """
        file_path = Path(path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        doc_type = detect_type(file_path)
        resolved_title = title or file_path.name
        source = str(file_path.resolve())

        logger.debug("Parsing file %s as %s", source, doc_type.value)

        content = self._read_file(file_path, doc_type)
        doc = Document(
            title=resolved_title,
            source=source,
            type=doc_type,
            content=content,
            metadata={
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "suffix": file_path.suffix,
                **(metadata or {}),
            },
        )
        if auto_chunk and content.strip():
            doc.chunks = self._chunker.chunk_document(doc)
        return doc

    def parse_text(
        self,
        text: str,
        *,
        title: str = "Untitled",
        source: str = "",
        doc_type: DocumentType = DocumentType.PLAIN_TEXT,
        metadata: dict[str, Any] | None = None,
        auto_chunk: bool = True,
    ) -> Document:
        """Parse raw text into a :class:`Document`.

        If ``doc_type`` is :attr:`DocumentType.MARKDOWN` or
        :attr:`DocumentType.HTML`, the text is converted to plain text
        before storage.
        """
        if doc_type is DocumentType.MARKDOWN:
            content = parse_markdown(text)
        elif doc_type is DocumentType.HTML:
            content = parse_html(text)
        else:
            content = text

        doc = Document(
            title=title,
            source=source,
            type=doc_type,
            content=content,
            metadata=metadata or {},
        )
        if auto_chunk and content.strip():
            doc.chunks = self._chunker.chunk_document(doc)
        return doc

    def parse_bytes(
        self,
        data: bytes,
        *,
        filename: str,
        title: str | None = None,
        metadata: dict[str, Any] | None = None,
        auto_chunk: bool = True,
    ) -> Document:
        """Parse raw bytes (e.g. from an upload) into a :class:`Document`.

        The document type is inferred from ``filename``.
        """
        doc_type = detect_type(filename)
        resolved_title = title or Path(filename).name

        if doc_type is DocumentType.PDF:
            content = _parse_pdf(data, filename)
        elif doc_type is DocumentType.WORD:
            content = _parse_word(data, filename)
        elif doc_type is DocumentType.EXCEL:
            content = _parse_excel(data, filename)
        elif doc_type is DocumentType.PPT:
            content = _parse_ppt(data, filename)
        elif doc_type is DocumentType.MARKDOWN:
            content = parse_markdown(data.decode("utf-8", errors="replace"))
        elif doc_type is DocumentType.HTML:
            content = parse_html(data.decode("utf-8", errors="replace"))
        else:
            content = data.decode("utf-8", errors="replace")

        doc = Document(
            title=resolved_title,
            source=filename,
            type=doc_type,
            content=content,
            metadata={
                "file_name": filename,
                "file_size": len(data),
                **(metadata or {}),
            },
        )
        if auto_chunk and content.strip():
            doc.chunks = self._chunker.chunk_document(doc)
        return doc

    # ------------------------------------------------------------------
    # Internal
    # ------------------------------------------------------------------

    def _read_file(self, path: Path, doc_type: DocumentType) -> str:
        """Read file content based on its type."""
        if doc_type is DocumentType.PDF:
            return _parse_pdf(path.read_bytes(), str(path))
        if doc_type is DocumentType.WORD:
            return _parse_word(path.read_bytes(), str(path))
        if doc_type is DocumentType.EXCEL:
            return _parse_excel(path.read_bytes(), str(path))
        if doc_type is DocumentType.PPT:
            return _parse_ppt(path.read_bytes(), str(path))
        if doc_type is DocumentType.MARKDOWN:
            return parse_markdown(path.read_text(encoding="utf-8", errors="replace"))
        if doc_type is DocumentType.HTML:
            return parse_html(path.read_text(encoding="utf-8", errors="replace"))
        # Default: plain text.
        return path.read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Document lifecycle manager
# ---------------------------------------------------------------------------


class DocumentLifecycleManager:
    """Manage document versioning, archival, and deletion.

    The manager holds documents in an in-memory registry. Every mutation
    that changes content creates a new :class:`DocumentVersion` snapshot
    so previous states can be restored. Archived documents are excluded
    from active queries but remain in the registry. Deleted documents are
    soft-deleted (status set to ``deleted``) and can be purged with
    :meth:`purge_deleted`.

    Example::

        >>> mgr = DocumentLifecycleManager()
        >>> doc = Document(title="README", content="Hello")
        >>> mgr.register(doc)
        >>> mgr.update_content(doc.id, "Hello world")
        >>> doc.version
        2
        >>> mgr.restore_version(doc.id, 1)
        >>> doc.content
        'Hello'
    """

    def __init__(self) -> None:
        self._documents: dict[str, Document] = {}

    # ------------------------------------------------------------------
    # Registration
    # ------------------------------------------------------------------

    def register(self, document: Document) -> Document:
        """Register a new document. Creates the first version snapshot."""
        if document.id in self._documents:
            raise ValueError(f"Document already registered: {document.id}")
        # Capture initial version.
        if not document.versions:
            document.versions.append(
                DocumentVersion(
                    version=1,
                    content=document.content,
                    content_hash=document.content_hash,
                    created_at=document.created_at,
                    metadata=dict(document.metadata),
                )
            )
        self._documents[document.id] = document
        logger.info("Registered document %s (%s)", document.id, document.title)
        return document

    def get(self, document_id: str) -> Document | None:
        """Return a document by ID, or None if not found."""
        return self._documents.get(document_id)

    def list_documents(
        self,
        *,
        status: DocumentStatus | None = None,
        doc_type: DocumentType | None = None,
    ) -> list[Document]:
        """List documents, optionally filtered by status or type."""
        result = list(self._documents.values())
        if status is not None:
            result = [d for d in result if d.status is status]
        if doc_type is not None:
            result = [d for d in result if d.type is doc_type]
        return result

    def __len__(self) -> int:
        return len(self._documents)

    def __contains__(self, document_id: str) -> bool:
        return document_id in self._documents

    def __iter__(self) -> Iterator[Document]:
        return iter(self._documents.values())

    # ------------------------------------------------------------------
    # Content updates & versioning
    # ------------------------------------------------------------------

    def update_content(
        self,
        document_id: str,
        content: str,
        *,
        chunker: TextChunker | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> Document:
        """Update a document's content, creating a new version snapshot.

        Args:
            document_id: ID of the document to update.
            content: New content text.
            chunker: Optional chunker to re-chunk the document. If None,
                no re-chunking is performed.
            metadata: Optional metadata to merge into the document.

        Raises:
            KeyError: If the document is not registered.
        """
        doc = self._require(document_id)
        new_version = doc.version + 1
        doc.versions.append(
            DocumentVersion(
                version=new_version,
                content=content,
                content_hash=_sha256_text(content),
                created_at=time.time(),
                metadata=dict(doc.metadata),
            )
        )
        doc.content = content
        doc.content_hash = _sha256_text(content)
        doc.version = new_version
        doc.updated_at = time.time()
        if metadata:
            doc.metadata.update(metadata)
        if chunker is not None:
            doc.chunks = chunker.chunk_document(doc)
        logger.info(
            "Updated document %s to version %d", document_id, new_version
        )
        return doc

    def update_metadata(
        self,
        document_id: str,
        metadata: dict[str, Any],
    ) -> Document:
        """Merge ``metadata`` into a document without creating a new version."""
        doc = self._require(document_id)
        doc.metadata.update(metadata)
        doc.updated_at = time.time()
        return doc

    def restore_version(self, document_id: str, version: int) -> Document:
        """Restore a document to a previous version's content.

        The restoration itself creates a new version (so the history is
        linear and never lost).

        Raises:
            KeyError: If the document or version is not found.
        """
        doc = self._require(document_id)
        snapshot = next(
            (v for v in doc.versions if v.version == version), None
        )
        if snapshot is None:
            raise KeyError(
                f"Version {version} not found for document {document_id}"
            )
        return self.update_content(document_id, snapshot.content)

    def get_version_history(self, document_id: str) -> list[DocumentVersion]:
        """Return the version history for a document."""
        doc = self._require(document_id)
        return list(doc.versions)

    # ------------------------------------------------------------------
    # Archival & deletion
    # ------------------------------------------------------------------

    def archive(self, document_id: str) -> Document:
        """Archive a document (exclude from active search)."""
        doc = self._require(document_id)
        doc.status = DocumentStatus.ARCHIVED
        doc.updated_at = time.time()
        logger.info("Archived document %s", document_id)
        return doc

    def unarchive(self, document_id: str) -> Document:
        """Restore an archived document to active status."""
        doc = self._require(document_id)
        doc.status = DocumentStatus.ACTIVE
        doc.updated_at = time.time()
        logger.info("Unarchived document %s", document_id)
        return doc

    def delete(self, document_id: str) -> Document:
        """Soft-delete a document."""
        doc = self._require(document_id)
        doc.status = DocumentStatus.DELETED
        doc.updated_at = time.time()
        logger.info("Soft-deleted document %s", document_id)
        return doc

    def purge_deleted(self) -> int:
        """Permanently remove all soft-deleted documents.

        Returns the number of documents purged.
        """
        to_purge = [
            doc_id
            for doc_id, doc in self._documents.items()
            if doc.status is DocumentStatus.DELETED
        ]
        for doc_id in to_purge:
            del self._documents[doc_id]
        if to_purge:
            logger.info("Purged %d deleted documents", len(to_purge))
        return len(to_purge)

    # ------------------------------------------------------------------
    # Internal
    # ------------------------------------------------------------------

    def _require(self, document_id: str) -> Document:
        doc = self._documents.get(document_id)
        if doc is None:
            raise KeyError(f"Document not found: {document_id}")
        return doc


__all__ = [
    "Chunk",
    "Document",
    "DocumentLifecycleManager",
    "DocumentParser",
    "DocumentStatus",
    "DocumentType",
    "DocumentVersion",
    "TextChunker",
    "detect_type",
    "format_timestamp",
    "parse_html",
    "parse_markdown",
]
